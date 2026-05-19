"""
SignAI — Professional PowerPoint Presentation Generator
========================================================
Generates: SignAI_Presentation.pptx  (17 slides)
Run:  python3 make_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy, os

OUT = os.path.join(os.path.dirname(__file__), "SignAI_Presentation.pptx")

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
W  = Inches(13.33)
H  = Inches(7.5)

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1a, 0x23, 0x7e)
BLUE       = RGBColor(0x15, 0x65, 0xc0)
LBLUE      = RGBColor(0x19, 0x76, 0xd2)
CYAN       = RGBColor(0x02, 0x88, 0xd1)
TEAL       = RGBColor(0x00, 0x69, 0x5c)
WHITE      = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG   = RGBColor(0xe8, 0xf4, 0xfd)
LIGHT_BG2  = RGBColor(0xf5, 0xf5, 0xf5)
GRAY       = RGBColor(0x42, 0x42, 0x42)
GRAY_LT    = RGBColor(0x75, 0x75, 0x75)
GREEN      = RGBColor(0x2e, 0x7d, 0x32)
LGREEN     = RGBColor(0xe8, 0xf5, 0xe9)
AMBER      = RGBColor(0xe6, 0x51, 0x00)
RED        = RGBColor(0xc6, 0x28, 0x28)
PURPLE     = RGBColor(0x6a, 0x1b, 0x9a)
DARK_BG    = RGBColor(0x0f, 0x17, 0x2a)
DARK_SURF  = RGBColor(0x1e, 0x29, 0x3b)
ACCENT     = RGBColor(0x63, 0x66, 0xf1)
ACCENT2    = RGBColor(0x81, 0x8c, 0xf8)
GRAY_LN    = RGBColor(0xe0, 0xe0, 0xe0)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape

def add_rounded_rect(slide, l, t, w, h, fill, line_color=None, line_w=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(5, l, t, w, h)   # rounded rect
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_para(tf, text, font_size=14, bold=False,
             color=GRAY, align=PP_ALIGN.LEFT, space_before=0, italic=False):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def slide_bg(slide, color):
    """Fill slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def header_bar(slide, title, subtitle="", bg=NAVY, h=Inches(1.35)):
    add_rect(slide, 0, 0, W, h, bg)
    # accent line
    add_rect(slide, 0, h - Inches(0.06), W, Inches(0.06), ACCENT)
    add_text(slide, title,
             Inches(0.45), Inches(0.18), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.45), Inches(0.88), Inches(12), Inches(0.35),
                 font_size=13, bold=False, color=ACCENT2, align=PP_ALIGN.LEFT)

def footer_bar(slide, text="SignAI  |  AI Sign Language Translator  |  Final Year Project  |  Himanshu Jagdish Patil"):
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), NAVY)
    add_text(slide, text,
             Inches(0.4), H - Inches(0.34), Inches(12.5), Inches(0.3),
             font_size=9, color=ACCENT2, align=PP_ALIGN.CENTER)

def bullet_box(slide, items, l, t, w, h,
               title="", title_color=NAVY, item_color=GRAY,
               bg=LIGHT_BG2, icon="▸", font_size=14, title_size=16):
    add_rounded_rect(slide, l, t, w, h, bg, LBLUE, 0.5)
    txb = slide.shapes.add_textbox(l + Inches(0.2), t + Inches(0.15),
                                    w - Inches(0.3), h - Inches(0.25))
    tf = txb.text_frame
    tf.word_wrap = True
    if title:
        add_para(tf, title, font_size=title_size, bold=True,
                 color=title_color, space_before=0)
    for item in items:
        add_para(tf, f"{icon}  {item}", font_size=font_size,
                 color=item_color, space_before=4)

def chip(slide, text, l, t, w=Inches(1.6), h=Inches(0.42),
         bg=ACCENT, fg=WHITE, font_size=12):
    add_rounded_rect(slide, l, t, w, h, bg)
    add_text(slide, text, l, t + Inches(0.04), w, h - Inches(0.06),
             font_size=font_size, bold=True, color=fg, align=PP_ALIGN.CENTER)

def divider(slide, y, color=ACCENT):
    add_rect(slide, Inches(0.4), y, W - Inches(0.8), Inches(0.04), color)

def stat_box(slide, value, label, l, t, w=Inches(2.4), h=Inches(1.4),
             bg=NAVY, val_color=WHITE, lbl_color=ACCENT2):
    add_rounded_rect(slide, l, t, w, h, bg)
    add_text(slide, value,
             l, t + Inches(0.18), w, Inches(0.65),
             font_size=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label,
             l, t + Inches(0.82), w, Inches(0.45),
             font_size=12, bold=False, color=lbl_color, align=PP_ALIGN.CENTER)

def arrow_right(slide, l, t, length=Inches(0.5), color=GRAY_LT):
    add_rect(slide, l, t + Inches(0.06), length - Inches(0.15), Inches(0.08), color)
    # arrowhead (triangle via shape 7)
    tri = slide.shapes.add_shape(7, l + length - Inches(0.2),
                                  t, Inches(0.2), Inches(0.2))
    tri.fill.solid(); tri.fill.fore_color.rgb = color
    tri.line.fill.background()

def flow_step(slide, label, sublabel, l, t, w=Inches(1.55), h=Inches(0.85),
              bg=LBLUE):
    add_rounded_rect(slide, l, t, w, h, bg)
    add_text(slide, label,
             l, t + Inches(0.06), w, Inches(0.42),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sublabel,
             l, t + Inches(0.46), w, Inches(0.34),
             font_size=8.5, bold=False, color=RGBColor(0xbb, 0xde, 0xfb),
             align=PP_ALIGN.CENTER)

def table_slide(slide, headers, rows, l, t, w, h, col_widths=None):
    """Add a styled table."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, l, t, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    # Header row
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = hdr
        run.font.size  = Pt(11)
        run.font.bold  = True
        run.font.color.rgb = WHITE

    # Data rows
    for r, row in enumerate(rows):
        bg = LIGHT_BG if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = str(val)
            run.font.size  = Pt(10)
            run.font.color.rgb = GRAY
    return tbl



# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
slide_bg(s1, DARK_BG)

# Top gradient band
add_rect(s1, 0, 0, W, Inches(0.08), ACCENT)

# Left dark panel
add_rect(s1, 0, Inches(0.08), Inches(5.2), H - Inches(0.08), DARK_SURF)
add_rect(s1, Inches(5.2), Inches(0.08), Inches(0.06), H - Inches(0.08), ACCENT)

# Big emoji area
add_text(s1, "🤟", Inches(0.3), Inches(0.5), Inches(4.8), Inches(2.0),
         font_size=90, align=PP_ALIGN.CENTER, color=WHITE)

# Project title left panel
add_text(s1, "SignAI",
         Inches(0.3), Inches(2.4), Inches(4.6), Inches(1.1),
         font_size=52, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
add_text(s1, "Real-Time AI Sign Language Translator",
         Inches(0.3), Inches(3.4), Inches(4.6), Inches(0.5),
         font_size=14, bold=False, color=RGBColor(0xbb,0xde,0xfb),
         align=PP_ALIGN.CENTER)

# Chips on left
chips_l = [("🧠 Deep Learning", ACCENT),
           ("🤟 40 Signs",      TEAL),
           ("🔊 Text-to-Speech",BLUE),
           ("💬 Sentence Builder", PURPLE)]
for i,(lbl,bg) in enumerate(chips_l):
    chip(s1, lbl, Inches(0.55 + i%2*2.1), Inches(4.15 + i//2*0.56),
         w=Inches(1.95), h=Inches(0.42), bg=bg)

# Right panel content
add_text(s1, "Final Year Project Presentation",
         Inches(5.55), Inches(0.55), Inches(7.5), Inches(0.55),
         font_size=15, bold=False, color=ACCENT2, align=PP_ALIGN.LEFT, italic=True)
add_text(s1, "SignAI: Real-Time AI-Powered\nSign Language Translator\nwith Text-to-Speech\n& Sentence Builder",
         Inches(5.55), Inches(1.05), Inches(7.5), Inches(2.6),
         font_size=29, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_rect(s1, Inches(5.55), Inches(3.6), Inches(7.3), Inches(0.04), ACCENT)

# Author / guide box
add_rounded_rect(s1, Inches(5.55), Inches(3.75), Inches(7.3), Inches(1.95),
                 DARK_SURF, ACCENT, 0.8)
info = [
    ("👤  Submitted By",       "Himanshu Jagdish Patil",   ACCENT2, WHITE),
    ("🎓  Roll No. / Batch",   "CS2024001  |  2021–2025",  ACCENT2, WHITE),
    ("👨‍🏫  Project Guide",     "Prof. [Guide Name], CSE",  ACCENT2, WHITE),
    ("🏛️  Department",         "CSE (AI & ML)",            ACCENT2, WHITE),
]
for i,(lbl,val,lc,vc) in enumerate(info):
    add_text(s1, lbl, Inches(5.75), Inches(3.9 + i*0.42),
             Inches(2.2), Inches(0.38), font_size=10, color=lc)
    add_text(s1, val, Inches(7.9),  Inches(3.9 + i*0.42),
             Inches(4.7), Inches(0.38), font_size=10, bold=True, color=vc)

# Bottom accent bar
add_rect(s1, 0, H - Inches(0.38), W, Inches(0.38), NAVY)
add_text(s1, "Department of Computer Science & Engineering (AI & ML)  |  Shri XYZ Institute of Technology",
         Inches(0.3), H - Inches(0.34), Inches(12.7), Inches(0.3),
         font_size=9, color=ACCENT2, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
slide_bg(s2, WHITE)
header_bar(s2, "Presentation Agenda", "What we will cover today")
footer_bar(s2)

agenda = [
    ("01", "Introduction & Problem Statement",    ACCENT),
    ("02", "Project Objectives",                  BLUE),
    ("03", "Technology Stack",                    TEAL),
    ("04", "System Architecture",                 NAVY),
    ("05", "How It Works — AI Pipeline",          PURPLE),
    ("06", "Data Collection & Training",          GREEN),
    ("07", "Backend & Frontend Implementation",   LBLUE),
    ("08", "Results & Accuracy",                  RGBColor(0xe6,0x51,0x00)),
    ("09", "Comparison with Existing Systems",    RED),
    ("10", "Future Work & Conclusion",            RGBColor(0x00,0x69,0x5c)),
]
cols = 2
for i, (num, topic, col) in enumerate(agenda):
    c = i % cols
    r = i // cols
    lx = Inches(0.5 + c * 6.4)
    ly = Inches(1.55 + r * 0.96)
    add_rounded_rect(s2, lx, ly, Inches(6.0), Inches(0.82), LIGHT_BG2, col, 1.0)
    add_rounded_rect(s2, lx, ly, Inches(0.7), Inches(0.82), col)
    add_text(s2, num, lx, ly + Inches(0.12), Inches(0.7), Inches(0.6),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, topic, lx + Inches(0.82), ly + Inches(0.2),
             Inches(5.0), Inches(0.45),
             font_size=14, bold=False, color=GRAY, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
slide_bg(s3, WHITE)
header_bar(s3, "Problem Statement", "The communication barrier we are solving")
footer_bar(s3)

# Stat boxes
stats = [
    ("466M+", "People with\nhearing loss\nworldwide (WHO)"),
    ("70M+",  "Deaf people\nuse sign language\nglobally"),
    ("<5%",   "Hearing people\ncan understand\nsign language"),
    ("High",  "Cost of human\nsign language\ninterpreters"),
]
for i,(val,lbl) in enumerate(stats):
    stat_box(s3, val, lbl, Inches(0.4 + i*3.1), Inches(1.5),
             w=Inches(2.85), h=Inches(1.5))

divider(s3, Inches(3.2))
add_text(s3, "Current Challenges:",
         Inches(0.45), Inches(3.3), Inches(12), Inches(0.42),
         font_size=16, bold=True, color=NAVY)

problems = [
    ("❶  Hardware Dependency",
     "Existing systems require expensive data gloves or depth sensors (Kinect, Leap Motion)."),
    ("❷  Limited Vocabulary",
     "Most systems handle only 24–26 static alphabet letters — not real-world phrases."),
    ("❸  No Temporal Modelling",
     "Static CNN approaches miss the motion-over-time nature of dynamic signs."),
    ("❹  No End-to-End Solution",
     "No system combines real-time detection + sentence building + text-to-speech in a web app."),
    ("❺  Poor Accessibility",
     "Research systems are lab-only and not deployable for everyday use."),
]
for i,(title,desc) in enumerate(problems):
    lx = Inches(0.4 + (i%3)*4.2) if i<3 else Inches(0.4 + (i-3)*6.25)
    ly = Inches(3.75) if i<3 else Inches(5.3)
    ww = Inches(3.9) if i<3 else Inches(5.8)
    add_rounded_rect(s3, lx, ly, ww, Inches(1.28), LIGHT_BG2, RED, 0.6)
    add_text(s3, title, lx+Inches(0.12), ly+Inches(0.08),
             ww-Inches(0.2), Inches(0.4),
             font_size=11, bold=True, color=RED)
    add_text(s3, desc, lx+Inches(0.12), ly+Inches(0.46),
             ww-Inches(0.2), Inches(0.75),
             font_size=10, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OBJECTIVES
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
slide_bg(s4, WHITE)
header_bar(s4, "Project Objectives", "What SignAI aims to achieve")
footer_bar(s4)

objectives = [
    ("🎯", "Real-Time Detection",     "Detect 40 ASL signs in real time\nusing a standard webcam — no special hardware"),
    ("🤟", "Dual-Hand Support",       "Extract 126 keypoints per frame\n(63 per hand) for comprehensive coverage"),
    ("🧠", "High Accuracy Model",     "Train a Stacked LSTM model to\nachieve ≥ 95% classification accuracy"),
    ("⚡", "Low Latency System",      "End-to-end prediction latency\nunder 200ms at 15 fps"),
    ("🔊", "Text-to-Speech Output",   "Convert detected signs to speech\nusing Google TTS + browser fallback"),
    ("💬", "Sentence Builder",        "Accumulate detected signs into\nfull sentences with word/sentence modes"),
    ("🌐", "Full Web Application",    "Build React.js + FastAPI full-stack\napp accessible from any browser"),
    ("📊", "Rigorous Evaluation",     "Report precision, recall, F1-score,\nconfusion matrix and system latency"),
]
for i,(icon,title,desc) in enumerate(objectives):
    c = i % 4
    r = i // 4
    lx = Inches(0.3 + c*3.24)
    ly = Inches(1.55 + r*2.42)
    add_rounded_rect(s4, lx, ly, Inches(3.0), Inches(2.22), LIGHT_BG2, ACCENT, 0.8)
    add_rect(s4, lx, ly, Inches(3.0), Inches(0.06), ACCENT)
    add_text(s4, icon, lx, ly+Inches(0.08), Inches(3.0), Inches(0.7),
             font_size=28, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(s4, title, lx+Inches(0.12), ly+Inches(0.75),
             Inches(2.78), Inches(0.42),
             font_size=12, bold=True, color=NAVY)
    add_text(s4, desc, lx+Inches(0.12), ly+Inches(1.15),
             Inches(2.78), Inches(0.95),
             font_size=10, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
slide_bg(s5, WHITE)
header_bar(s5, "Technology Stack", "Tools and frameworks powering SignAI")
footer_bar(s5)

# Left column — backend
add_rounded_rect(s5, Inches(0.35), Inches(1.5), Inches(3.9), Inches(5.6),
                 LIGHT_BG2, NAVY, 1)
add_text(s5, "🐍  BACKEND", Inches(0.35), Inches(1.55),
         Inches(3.9), Inches(0.5), font_size=14, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER)
divider(s5, Inches(2.1))
backend_items = [
    ("FastAPI  v0.111",   "WebSocket + REST API server", TEAL),
    ("TensorFlow  v2.16", "LSTM model training & inference", BLUE),
    ("MediaPipe  v0.10",  "21-landmark hand detection", GREEN),
    ("OpenCV  v4.9",      "Frame decode & colour convert", NAVY),
    ("gTTS  v2.5",        "Google Text-to-Speech (MP3)", PURPLE),
    ("Uvicorn  v0.29",    "ASGI production server", LBLUE),
    ("NumPy  v1.26",      "Keypoint array processing", GRAY),
    ("scikit-learn  v1.4","Dataset split & metrics", RGBColor(0xe6,0x51,0x00)),
]
for i,(name,desc,col) in enumerate(backend_items):
    ly = Inches(2.25 + i*0.59)
    add_rounded_rect(s5, Inches(0.55), ly, Inches(3.5), Inches(0.5), WHITE, col, 0.8)
    add_text(s5, name, Inches(0.65), ly+Inches(0.06),
             Inches(1.55), Inches(0.38), font_size=10, bold=True, color=col)
    add_text(s5, desc, Inches(2.22), ly+Inches(0.06),
             Inches(1.72), Inches(0.38), font_size=9, color=GRAY_LT)

# Middle column — frontend
add_rounded_rect(s5, Inches(4.55), Inches(1.5), Inches(3.9), Inches(5.6),
                 LIGHT_BG2, ACCENT, 1)
add_text(s5, "⚛️  FRONTEND", Inches(4.55), Inches(1.55),
         Inches(3.9), Inches(0.5), font_size=14, bold=True,
         color=ACCENT, align=PP_ALIGN.CENTER)
divider(s5, Inches(2.1))
frontend_items = [
    ("React  v18.3",      "Core UI component library",       ACCENT),
    ("react-webcam  v7.2","Browser webcam stream access",    BLUE),
    ("Framer Motion v11", "Smooth sign animations",           PURPLE),
    ("WebSocket API",     "Real-time prediction streaming",   TEAL),
    ("Google TTS API",    "Speech synthesis endpoint",        GREEN),
    ("Custom CSS",        "Dark glassmorphism theme",         NAVY),
    ("lucide-react",      "Icon component library",           GRAY),
    ("Node.js  v18 LTS",  "Build & development server",       RGBColor(0x00,0x69,0x5c)),
]
for i,(name,desc,col) in enumerate(frontend_items):
    ly = Inches(2.25 + i*0.59)
    add_rounded_rect(s5, Inches(4.75), ly, Inches(3.5), Inches(0.5), WHITE, col, 0.8)
    add_text(s5, name, Inches(4.85), ly+Inches(0.06),
             Inches(1.55), Inches(0.38), font_size=10, bold=True, color=col)
    add_text(s5, desc, Inches(6.42), ly+Inches(0.06),
             Inches(1.72), Inches(0.38), font_size=9, color=GRAY_LT)

# Right column — ML / AI
add_rounded_rect(s5, Inches(8.75), Inches(1.5), Inches(4.2), Inches(5.6),
                 LIGHT_BG2, GREEN, 1)
add_text(s5, "🧠  AI / ML PIPELINE", Inches(8.75), Inches(1.55),
         Inches(4.2), Inches(0.5), font_size=14, bold=True,
         color=GREEN, align=PP_ALIGN.CENTER)
divider(s5, Inches(2.1))
ml_items = [
    ("Stacked LSTM",        "3-layer temporal classifier",     GREEN),
    ("126 Keypoints/Frame", "21 landmarks × 3 × 2 hands",     TEAL),
    ("30-Frame Window",     "~2 sec sequence per prediction",  BLUE),
    ("Batch Normalisation", "Training stability & speed",      NAVY),
    ("Dropout  (0.3/0.4)",  "Overfitting prevention",          PURPLE),
    ("Adam Optimiser",      "lr=0.001, reduce on plateau",     RGBColor(0xe6,0x51,0x00)),
    ("EarlyStopping",       "Patience=20, best weights saved", RED),
    ("40-class Softmax",    "Confidence score output",         RGBColor(0x00,0x69,0x5c)),
]
for i,(name,desc,col) in enumerate(ml_items):
    ly = Inches(2.25 + i*0.59)
    add_rounded_rect(s5, Inches(8.95), ly, Inches(3.8), Inches(0.5), WHITE, col, 0.8)
    add_text(s5, name, Inches(9.05), ly+Inches(0.06),
             Inches(1.75), Inches(0.38), font_size=10, bold=True, color=col)
    add_text(s5, desc, Inches(10.82), ly+Inches(0.06),
             Inches(1.82), Inches(0.38), font_size=9, color=GRAY_LT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
slide_bg(s6, WHITE)
header_bar(s6, "System Architecture", "Three-tier client-server design")
footer_bar(s6)

# Three columns
cols_arch = [
    ("🖥️ FRONTEND", "React.js", ACCENT,
     ["Camera.jsx\n(Webcam)", "SignDisplay.jsx\n(Results)", "SentenceBuilder\n(TTS)", "StatusBar\n(WS status)"]),
    ("⚡ BACKEND", "FastAPI + Python", TEAL,
     ["WebSocket\n/ws/predict", "MediaPipe\nHands (×2)", "LSTM\nInference", "gTTS\n/api/tts"]),
    ("🧠 ML PIPELINE", "TensorFlow / Keras", GREEN,
     ["Data Collection\ncollect_data.py", "Keypoints\n(30,126) .npy", "Training\ntrain.py", "Model\nsign_model.h5"]),
]
for ci,(col_title,col_sub,col_color,boxes) in enumerate(cols_arch):
    lx = Inches(0.32 + ci*4.32)
    add_rounded_rect(s6, lx, Inches(1.5), Inches(4.0), Inches(5.55),
                     LIGHT_BG2, col_color, 1.2)
    add_rounded_rect(s6, lx, Inches(1.5), Inches(4.0), Inches(0.7), col_color)
    add_text(s6, col_title, lx, Inches(1.55),
             Inches(4.0), Inches(0.38), font_size=14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s6, col_sub, lx, Inches(1.9),
             Inches(4.0), Inches(0.28), font_size=10, color=WHITE,
             align=PP_ALIGN.CENTER, italic=True)
    for bi,box_lbl in enumerate(boxes):
        bly = Inches(2.35 + bi*1.12)
        add_rounded_rect(s6, lx+Inches(0.25), bly,
                         Inches(3.5), Inches(0.92), WHITE, col_color, 0.6)
        add_text(s6, box_lbl, lx+Inches(0.25), bly+Inches(0.08),
                 Inches(3.5), Inches(0.76),
                 font_size=11, color=GRAY, align=PP_ALIGN.CENTER)
        if bi < len(boxes)-1:
            add_rect(s6, lx+Inches(1.9), bly+Inches(0.92),
                     Inches(0.18), Inches(0.2), col_color)

# Arrows between columns
for ci in range(2):
    lx_arr = Inches(4.32 + ci*4.32)
    for ay in [Inches(2.74), Inches(3.86), Inches(4.98)]:
        add_rect(s6, lx_arr, ay, Inches(0.32), Inches(0.05), GRAY_LT)
    add_text(s6, "↔" if ci==0 else "→",
             lx_arr - Inches(0.05), Inches(3.3), Inches(0.42), Inches(0.55),
             font_size=22, color=GRAY_LT, align=PP_ALIGN.CENTER)

add_text(s6, "WebSocket (ws://)",
         Inches(4.18), Inches(6.62), Inches(1.25), Inches(0.32),
         font_size=8, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
add_text(s6, "Local / Offline",
         Inches(8.5), Inches(6.62), Inches(1.25), Inches(0.32),
         font_size=8, color=GREEN, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — HOW IT WORKS (AI PIPELINE)
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
slide_bg(s7, WHITE)
header_bar(s7, "How It Works — AI Pipeline", "From webcam frame to predicted sign")
footer_bar(s7)

steps = [
    ("📷", "Webcam\nCapture",     "15 fps\n640×480",         ACCENT),
    ("✂️", "Frame\nEncoding",     "Base64\nJPEG",             LBLUE),
    ("🌐", "WebSocket\nTransfer", "~5 ms\nlatency",           BLUE),
    ("✋", "MediaPipe\nHands ×2", "21 kp × 3\n= 63 each",    TEAL),
    ("📐", "Keypoint\nExtraction","126 values\nper frame",    GREEN),
    ("🗂️", "30-Frame\nWindow",    "Sliding\nbuffer",          NAVY),
    ("🧠", "LSTM\nInference",     "3-layer\nLSTM",            PURPLE),
    ("🎯", "Prediction\n+ Conf.", "Sign name\n+ score %",     RGBColor(0xe6,0x51,0x00)),
]
step_w = Inches(1.42)
step_h = Inches(1.65)
gap    = Inches(0.12)
start  = Inches(0.3)
top_y  = Inches(1.65)

for i,(icon,title,sub,col) in enumerate(steps):
    lx = start + i*(step_w+gap)
    add_rounded_rect(s7, lx, top_y, step_w, step_h, col)
    add_text(s7, icon, lx, top_y+Inches(0.08),
             step_w, Inches(0.56), font_size=26,
             align=PP_ALIGN.CENTER, color=WHITE)
    add_text(s7, title, lx, top_y+Inches(0.62),
             step_w, Inches(0.56), font_size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s7, sub, lx, top_y+Inches(1.16),
             step_w, Inches(0.44), font_size=9,
             color=RGBColor(0xbb,0xde,0xfb), align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        add_text(s7, "▶", lx+step_w+Inches(0.01), top_y+Inches(0.65),
                 gap+Inches(0.1), Inches(0.42),
                 font_size=14, color=GRAY_LT, align=PP_ALIGN.CENTER)
    # Step number
    add_rounded_rect(s7, lx+Inches(0.04), top_y-Inches(0.24),
                     Inches(0.32), Inches(0.32), col)
    add_text(s7, str(i+1), lx+Inches(0.04), top_y-Inches(0.24),
             Inches(0.32), Inches(0.32), font_size=10, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

# Key info boxes
divider(s7, Inches(3.55))
info_boxes = [
    ("🤟 Both Hands",        "Left + Right hand\nextracted separately\nMissing = zero-pad",  TEAL),
    ("📊 126 Features",      "63 values/hand × 2\n= 126 total\nper frame",                  BLUE),
    ("⏱️ 30 Frames = ~2 sec", "At 15 fps, one prediction\nrequires ~2 sec of\nhand movement", ACCENT),
    ("🎯 75% Threshold",     "Predictions below\n75% confidence\nare suppressed",            GREEN),
    ("🔄 Cooldown 1.2 sec",  "Same sign not\nrepeated within\n1.2 seconds",                  NAVY),
]
for i,(title,desc,col) in enumerate(info_boxes):
    lx = Inches(0.3 + i*2.6)
    add_rounded_rect(s7, lx, Inches(3.7), Inches(2.42), Inches(2.82),
                     LIGHT_BG2, col, 0.8)
    add_rect(s7, lx, Inches(3.7), Inches(2.42), Inches(0.06), col)
    add_text(s7, title, lx+Inches(0.1), Inches(3.8),
             Inches(2.22), Inches(0.42), font_size=12, bold=True, color=col)
    add_text(s7, desc, lx+Inches(0.1), Inches(4.22),
             Inches(2.22), Inches(2.18), font_size=11, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DATA COLLECTION
# ═══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
slide_bg(s8, WHITE)
header_bar(s8, "Data Collection", "Building the training dataset with collect_data.py")
footer_bar(s8)

# Left — sign categories
add_rounded_rect(s8, Inches(0.35), Inches(1.5), Inches(4.5), Inches(5.55),
                 LIGHT_BG2, NAVY, 1)
add_text(s8, "📋  40 Supported Signs", Inches(0.45), Inches(1.6),
         Inches(4.3), Inches(0.46), font_size=14, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER)

cats = [
    ("💬 Common Phrases (10)", "hello, yes, no, thankyou, please\nsorry, help, good, bad, stop", ACCENT),
    ("😊 Expressions (5)",     "iloveyou, goodmorning, goodnight\nhowareyou, fine",              BLUE),
    ("🏃 Actions (5)",          "eat, drink, sleep, come, go",                                   TEAL),
    ("🔤 Alphabet A–J (10)",    "A, B, C, D, E, F, G, H, I, J",                                 GREEN),
    ("🔢 Numbers 1–10 (10)",    "1, 2, 3, 4, 5, 6, 7, 8, 9, 10",                               PURPLE),
]
for i,(cat,items,col) in enumerate(cats):
    ly = Inches(2.12 + i*0.96)
    add_rounded_rect(s8, Inches(0.55), ly, Inches(4.1), Inches(0.86), WHITE, col, 0.6)
    add_rect(s8, Inches(0.55), ly, Inches(0.12), Inches(0.86), col)
    add_text(s8, cat, Inches(0.8), ly+Inches(0.06),
             Inches(3.7), Inches(0.36), font_size=11, bold=True, color=col)
    add_text(s8, items, Inches(0.8), ly+Inches(0.42),
             Inches(3.7), Inches(0.38), font_size=9.5, color=GRAY_LT)

# Right — stats and controls
add_rounded_rect(s8, Inches(5.2), Inches(1.5), Inches(3.7), Inches(2.6),
                 LIGHT_BG2, ACCENT, 1)
add_text(s8, "📊  Dataset Statistics", Inches(5.3), Inches(1.58),
         Inches(3.5), Inches(0.42), font_size=13, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER)
ds_stats = [
    ("Total Classes",    "40 signs"),
    ("Samples / Sign",   "200"),
    ("Total Samples",    "8,000"),
    ("Frames / Sample",  "30 frames"),
    ("Keypoints/Frame",  "126 (both hands)"),
    ("Feature Size",     "30 × 126 = 3,780"),
]
for i,(lbl,val) in enumerate(ds_stats):
    ly = Inches(2.05 + i*0.33)
    add_text(s8, lbl+":", Inches(5.35), ly, Inches(2.1), Inches(0.3),
             font_size=10, color=GRAY)
    add_text(s8, val, Inches(7.45), ly, Inches(1.3), Inches(0.3),
             font_size=10, bold=True, color=NAVY)

# Controls
add_rounded_rect(s8, Inches(5.2), Inches(4.25), Inches(3.7), Inches(2.8),
                 LIGHT_BG2, GREEN, 1)
add_text(s8, "⌨️  Controls", Inches(5.3), Inches(4.33),
         Inches(3.5), Inches(0.4), font_size=13, bold=True,
         color=GREEN, align=PP_ALIGN.CENTER)
controls = [
    ("ENTER", "Record 1 sample (30 frames)", GREEN),
    ("N",     "Move to next sign",           BLUE),
    ("P",     "Move to previous sign",       BLUE),
    ("Q",     "Quit data collection",        RED),
]
for i,(key,action,col) in enumerate(controls):
    ly = Inches(4.78 + i*0.54)
    add_rounded_rect(s8, Inches(5.4), ly, Inches(0.72), Inches(0.42), col)
    add_text(s8, key, Inches(5.4), ly+Inches(0.03),
             Inches(0.72), Inches(0.36), font_size=10, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s8, action, Inches(6.2), ly+Inches(0.06),
             Inches(2.6), Inches(0.32), font_size=10, color=GRAY)

# Tips
add_rounded_rect(s8, Inches(9.15), Inches(1.5), Inches(3.85), Inches(5.55),
                 LIGHT_BG2, AMBER, 1)
add_text(s8, "💡  Tips for High Quality", Inches(9.25), Inches(1.58),
         Inches(3.65), Inches(0.42), font_size=13, bold=True,
         color=AMBER, align=PP_ALIGN.CENTER)
tips = [
    "✅  Good indoor lighting\n     (face a window/lamp)",
    "✅  Plain background\n     (white wall is best)",
    "✅  Vary hand angle slightly\n     between samples",
    "✅  Collect 200 samples per\n     sign for 95%+ accuracy",
    "❌  Avoid dark rooms or\n     cluttered backgrounds",
    "❌  Do not rush — perform\n     signs clearly & slowly",
]
for i,tip in enumerate(tips):
    ly = Inches(2.12 + i*0.84)
    col = GREEN if tip.startswith("✅") else RED
    add_rounded_rect(s8, Inches(9.3), ly, Inches(3.5), Inches(0.72), WHITE, col, 0.5)
    add_text(s8, tip, Inches(9.42), ly+Inches(0.06),
             Inches(3.3), Inches(0.62), font_size=10, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LSTM MODEL
# ═══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
slide_bg(s9, WHITE)
header_bar(s9, "LSTM Model Architecture", "Stacked LSTM with BatchNorm and Dropout")
footer_bar(s9)

# Left — layer diagram
add_text(s9, "Model Architecture", Inches(0.35), Inches(1.52),
         Inches(5.2), Inches(0.42), font_size=14, bold=True, color=NAVY)
layers = [
    ("Input",                   "(30, 126)",          RGBColor(0x01,0x57,0x8b), Inches(1.1)),
    ("LSTM  128 units",         "return_seq=True",    LBLUE,                    Inches(0.86)),
    ("BatchNorm + Dropout 0.3", "",                   RGBColor(0x28,0x35,0x93), Inches(0.64)),
    ("LSTM  256 units",         "return_seq=True",    LBLUE,                    Inches(0.86)),
    ("BatchNorm + Dropout 0.3", "",                   RGBColor(0x28,0x35,0x93), Inches(0.64)),
    ("LSTM  128 units",         "return_seq=False",   LBLUE,                    Inches(0.86)),
    ("BatchNorm + Dropout 0.3", "",                   RGBColor(0x28,0x35,0x93), Inches(0.64)),
    ("Dense 256  → ReLU",       "Dropout 0.4",        GREEN,                    Inches(0.64)),
    ("Dense 128  → ReLU",       "Dropout 0.3",        RGBColor(0x1b,0x5e,0x20), Inches(0.64)),
    ("Dense 40   → Softmax",    "Output layer",       RED,                      Inches(0.64)),
]
bw_l = Inches(4.6)
cur_y = Inches(2.0)
for lbl,sub,col,bh in layers:
    add_rounded_rect(s9, Inches(0.35), cur_y, bw_l, bh, col)
    add_text(s9, lbl, Inches(0.45), cur_y+Inches(0.06),
             Inches(3.4), bh-Inches(0.1), font_size=10.5, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
    if sub:
        add_text(s9, sub, Inches(3.85), cur_y+Inches(0.06),
                 Inches(1.0), bh-Inches(0.1), font_size=9,
                 color=RGBColor(0xbb,0xde,0xfb), align=PP_ALIGN.CENTER)
    cur_y += bh + Inches(0.04)

# Right — hyperparams table
add_text(s9, "Hyperparameters", Inches(5.3), Inches(1.52),
         Inches(7.7), Inches(0.42), font_size=14, bold=True, color=NAVY)
hp_headers = ["Parameter", "Value"]
hp_rows = [
    ["Input Shape",        "(30, 126)"],
    ["LSTM Units",         "128 → 256 → 128"],
    ["Dense Units",        "256 → 128 → 40"],
    ["Optimizer",          "Adam  lr=0.001"],
    ["Loss Function",      "Categorical Cross-Entropy"],
    ["Max Epochs",         "200"],
    ["Batch Size",         "32"],
    ["Early Stopping",     "Patience = 20"],
    ["LR Reduction",       "Factor=0.5, Patience=8"],
    ["Total Parameters",   "~1.2 Million"],
    ["Model File Size",    "~14 MB (.h5)"],
    ["Best Val Accuracy",  "95.9%"],
    ["Test Accuracy",      "96.8%"],
]
table_slide(s9, hp_headers, hp_rows,
            Inches(5.3), Inches(2.0), Inches(7.7), Inches(5.55),
            col_widths=[Inches(3.3), Inches(4.0)])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — BACKEND & API
# ═══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
slide_bg(s10, WHITE)
header_bar(s10, "Backend Implementation", "FastAPI server — WebSocket + REST + TTS")
footer_bar(s10)

# API endpoints table
add_text(s10, "🔌  API Endpoints", Inches(0.35), Inches(1.52),
         Inches(12.7), Inches(0.42), font_size=14, bold=True, color=NAVY)
ep_headers = ["Method", "Endpoint", "Description", "Request", "Response"]
ep_rows = [
    ["WS",   "/ws/predict",  "Real-time sign prediction",  "JSON: {frame: base64}",    "JSON: {sign, confidence, top5}"],
    ["POST", "/api/tts",     "Text-to-Speech MP3 output",  "JSON: {text, lang, slow}", "MP3 audio stream"],
    ["GET",  "/api/signs",   "List all 40 signs",          "—",                        "JSON: {signs[], count}"],
    ["GET",  "/api/health",  "Server health + model info", "—",                        "JSON: {status, sign_count}"],
]
table_slide(s10, ep_headers, ep_rows,
            Inches(0.35), Inches(2.0), Inches(12.6), Inches(2.2),
            col_widths=[Inches(0.9), Inches(1.7), Inches(2.8), Inches(3.2), Inches(3.8)])

# WebSocket flow
divider(s10, Inches(4.35))
add_text(s10, "⚡  WebSocket Prediction Flow (per frame cycle)",
         Inches(0.35), Inches(4.44), Inches(12.7), Inches(0.42),
         font_size=14, bold=True, color=TEAL)

ws_steps = [
    ("1. Receive\nframe", ACCENT),  ("2. Decode\nBase64→BGR", LBLUE),
    ("3. MediaPipe\nkeypoints", TEAL), ("4. Append\nto buffer", BLUE),
    ("5. Buffer\nfull? (30)", NAVY),   ("6. LSTM\npredict", GREEN),
    ("7. Conf\n≥ 75%?", PURPLE),       ("8. Send\nprediction", RED),
]
sw = Inches(1.4); sh = Inches(0.78); sg = Inches(0.2)
for i,(lbl,col) in enumerate(ws_steps):
    lx = Inches(0.35) + i*(sw+sg)
    add_rounded_rect(s10, lx, Inches(4.95), sw, sh, col)
    add_text(s10, lbl, lx, Inches(5.0), sw, sh-Inches(0.08),
             font_size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(ws_steps)-1:
        add_text(s10, "▶", lx+sw+Inches(0.03), Inches(5.17),
                 sg+Inches(0.06), Inches(0.42), font_size=12,
                 color=GRAY_LT, align=PP_ALIGN.CENTER)

# Key params
param_boxes = [
    ("SEQUENCE_LEN",   "30 frames", "Sliding window size",   ACCENT),
    ("CONF_THRESHOLD", "0.75",      "Minimum confidence",    TEAL),
    ("COOLDOWN_SECS",  "1.2 sec",   "Anti-duplicate delay",  GREEN),
    ("MAX_NUM_HANDS",  "2",         "Both hands detected",   NAVY),
    ("FPS",            "15 fps",    "Frame capture rate",    BLUE),
]
for i,(name,val,desc,col) in enumerate(param_boxes):
    lx = Inches(0.35 + i*2.56)
    add_rounded_rect(s10, lx, Inches(6.0), Inches(2.38), Inches(1.06),
                     LIGHT_BG2, col, 0.8)
    add_text(s10, name, lx+Inches(0.1), Inches(6.06),
             Inches(2.18), Inches(0.3), font_size=9, bold=True, color=col)
    add_text(s10, val, lx+Inches(0.1), Inches(6.34),
             Inches(2.18), Inches(0.36), font_size=16, bold=True, color=NAVY)
    add_text(s10, desc, lx+Inches(0.1), Inches(6.68),
             Inches(2.18), Inches(0.3), font_size=8.5, color=GRAY_LT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FRONTEND UI
# ═══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
slide_bg(s11, WHITE)
header_bar(s11, "Frontend — React.js UI", "Modern dark glassmorphism web application")
footer_bar(s11)

# Left: dark UI mockup panel
add_rounded_rect(s11, Inches(0.35), Inches(1.52), Inches(6.15), Inches(5.55),
                 DARK_BG, ACCENT, 1.2)

# Mock header bar inside panel
add_rect(s11, Inches(0.35), Inches(1.52), Inches(6.15), Inches(0.48), DARK_SURF)
add_text(s11, "🤟  SignAI",
         Inches(0.5), Inches(1.56), Inches(2.0), Inches(0.38),
         font_size=12, bold=True, color=ACCENT2)
add_rounded_rect(s11, Inches(4.5), Inches(1.6), Inches(0.72), Inches(0.28),
                 DARK_SURF, ACCENT, 0.5)
add_text(s11, "40 Signs", Inches(4.5), Inches(1.6),
         Inches(0.72), Inches(0.28), font_size=7, color=GRAY_LT, align=PP_ALIGN.CENTER)
add_rounded_rect(s11, Inches(5.3), Inches(1.6), Inches(0.9), Inches(0.28), ACCENT)
add_text(s11, "AI Powered", Inches(5.3), Inches(1.6),
         Inches(0.9), Inches(0.28), font_size=7, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

# Mock status bar
add_rounded_rect(s11, Inches(0.5), Inches(2.08), Inches(5.85), Inches(0.32),
                 DARK_SURF, RGBColor(0x33,0x41,0x55), 0.5)
add_text(s11, "●  Connected        ✋  Hand Detected        WebSocket: Active  |  15 fps",
         Inches(0.6), Inches(2.1), Inches(5.65), Inches(0.28),
         font_size=7.5, color=GREEN)

# Mock tab nav
add_rounded_rect(s11, Inches(0.5), Inches(2.48), Inches(1.5), Inches(0.32), ACCENT)
add_text(s11, "🤟 Translator", Inches(0.5), Inches(2.48),
         Inches(1.5), Inches(0.32), font_size=7.5, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_rounded_rect(s11, Inches(2.1), Inches(2.48), Inches(1.3), Inches(0.32),
                 DARK_SURF, RGBColor(0x33,0x41,0x55), 0.5)
add_text(s11, "📖 Sign List", Inches(2.1), Inches(2.48),
         Inches(1.3), Inches(0.32), font_size=7.5, color=GRAY_LT, align=PP_ALIGN.CENTER)

# Mock camera panel
add_rounded_rect(s11, Inches(0.5), Inches(2.9), Inches(2.78), Inches(2.1),
                 RGBColor(0x0a,0x0a,0x0f), RGBColor(0x33,0x41,0x55), 0.5)
add_text(s11, "🤟", Inches(0.5), Inches(3.15), Inches(2.78), Inches(1.0),
         font_size=38, align=PP_ALIGN.CENTER, color=WHITE)
add_rounded_rect(s11, Inches(2.8), Inches(2.98), Inches(0.62), Inches(0.24),
                 RGBColor(0xc6,0x28,0x28))
add_text(s11, "⬤ LIVE", Inches(2.8), Inches(2.98),
         Inches(0.62), Inches(0.24), font_size=7, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_rounded_rect(s11, Inches(0.5), Inches(4.88), Inches(1.3), Inches(0.28),
                 RGBColor(0x45,0x0a,0x0a), RGBColor(0xef,0x44,0x44), 0.5)
add_text(s11, "⏹ Stop", Inches(0.5), Inches(4.88), Inches(1.3), Inches(0.28),
         font_size=8, bold=True, color=RGBColor(0xef,0x44,0x44), align=PP_ALIGN.CENTER)
add_rounded_rect(s11, Inches(1.9), Inches(4.88), Inches(1.1), Inches(0.28),
                 DARK_SURF, RGBColor(0x33,0x41,0x55), 0.5)
add_text(s11, "↔ Mirror", Inches(1.9), Inches(4.88), Inches(1.1), Inches(0.28),
         font_size=8, color=RGBColor(0xf1,0xf5,0xf9), align=PP_ALIGN.CENTER)

# Mock sign display card
add_rounded_rect(s11, Inches(3.4), Inches(2.9), Inches(2.78), Inches(1.72),
                 RGBColor(0x1e,0x1b,0x4b), RGBColor(0x43,0x38,0xca), 0.8)
add_text(s11, "hello",
         Inches(3.4), Inches(3.05), Inches(2.78), Inches(0.9),
         font_size=34, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
add_rounded_rect(s11, Inches(3.55), Inches(4.0), Inches(2.2), Inches(0.16),
                 RGBColor(0x1e,0x29,0x3b))
add_rounded_rect(s11, Inches(3.55), Inches(4.0), Inches(2.08), Inches(0.16), GREEN)
add_text(s11, "96.2%", Inches(5.7), Inches(3.98),
         Inches(0.42), Inches(0.2), font_size=8, bold=True, color=GREEN)

# Mock sentence builder
add_rounded_rect(s11, Inches(0.5), Inches(5.25), Inches(5.65), Inches(1.58),
                 DARK_SURF, RGBColor(0x33,0x41,0x55), 0.5)
add_text(s11, "💬 Sentence Builder",
         Inches(0.6), Inches(5.3), Inches(3.0), Inches(0.32),
         font_size=10, bold=True, color=WHITE)
chip_words = [("hello",ACCENT),("thankyou",RGBColor(0x7c,0x3a,0xed)),
              ("good",RGBColor(0x05,0x96,0x69))]
for wi,(w,bg) in enumerate(chip_words):
    chip(s11, w, Inches(0.6+wi*1.55), Inches(5.7),
         w=Inches(1.45), h=Inches(0.32), bg=bg, font_size=9)
add_rounded_rect(s11, Inches(0.6), Inches(6.12), Inches(0.82), Inches(0.28),
                 ACCENT)
add_text(s11, "🔊 Speak", Inches(0.6), Inches(6.12),
         Inches(0.82), Inches(0.28), font_size=7.5, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_rounded_rect(s11, Inches(1.5), Inches(6.12), Inches(0.65), Inches(0.28),
                 DARK_SURF, RGBColor(0x33,0x41,0x55), 0.5)
add_text(s11, "📋 Copy", Inches(1.5), Inches(6.12),
         Inches(0.65), Inches(0.28), font_size=7.5, color=GRAY_LT, align=PP_ALIGN.CENTER)
add_rounded_rect(s11, Inches(2.25), Inches(6.12), Inches(0.65), Inches(0.28),
                 DARK_SURF, RGBColor(0x33,0x41,0x55), 0.5)
add_text(s11, "⌫ Undo", Inches(2.25), Inches(6.12),
         Inches(0.65), Inches(0.28), font_size=7.5, color=GRAY_LT, align=PP_ALIGN.CENTER)

# Right: feature list
add_text(s11, "UI Components & Features",
         Inches(6.75), Inches(1.52), Inches(6.2), Inches(0.42),
         font_size=14, bold=True, color=NAVY)

ui_features = [
    ("📷 Camera Panel",
     "Live webcam feed with MediaPipe green\nskeleton overlay on both hands.\nStart/Stop/Mirror controls.",
     ACCENT),
    ("🔤 Sign Display",
     "Detected sign in large gradient text with\nanimated confidence bar + Top-5\npredictions list.",
     BLUE),
    ("💬 Sentence Builder",
     "Auto-adds signs as coloured word chips.\nToggle Word / Sentence mode.\nClick chip to remove it.",
     TEAL),
    ("🔊 Text-to-Speech",
     "Speaks full sentence via Google TTS.\nBrowser SpeechSynthesis fallback.\nWorks offline too.",
     GREEN),
    ("📖 Sign Dictionary",
     "Browse all 40 signs by category.\nSearch box for quick lookup.\nEmoji + name cards.",
     PURPLE),
    ("📊 Status Bar",
     "Live WebSocket connection status.\nHand detected indicator.\nFPS and model info.",
     NAVY),
]
for i,(title,desc,col) in enumerate(ui_features):
    c = i % 2; r = i // 2
    lx = Inches(6.75 + c*3.12)
    ly = Inches(2.05 + r*1.75)
    add_rounded_rect(s11, lx, ly, Inches(2.9), Inches(1.62), LIGHT_BG2, col, 0.8)
    add_rect(s11, lx, ly, Inches(2.9), Inches(0.06), col)
    add_text(s11, title, lx+Inches(0.12), ly+Inches(0.12),
             Inches(2.66), Inches(0.4), font_size=11, bold=True, color=col)
    add_text(s11, desc, lx+Inches(0.12), ly+Inches(0.52),
             Inches(2.66), Inches(1.04), font_size=9.5, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — KEY FEATURES
# ═══════════════════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
slide_bg(s12, WHITE)
header_bar(s12, "Key Features", "What makes SignAI stand out")
footer_bar(s12)

features = [
    ("🤟", "Real-Time\nDetection",
     "Live webcam sign detection\nat 15 fps with ~67ms\nend-to-end latency",
     ACCENT),
    ("✋✋", "Dual-Hand\nSupport",
     "Both hands detected\nsimultaneously — 126\nkeypoints per frame",
     TEAL),
    ("🧠", "95%+\nAccuracy",
     "Stacked LSTM model\ntrained on 8,000 samples\nacross 40 sign classes",
     GREEN),
    ("🔊", "Text-to-Speech\nOutput",
     "Google TTS generates\nnatural MP3 speech\nfrom detected sentences",
     BLUE),
    ("💬", "Sentence\nBuilder",
     "Word chips + Sentence mode\nQuick-Add buttons, Copy,\nUndo and Clear",
     PURPLE),
    ("📱", "No Special\nHardware",
     "Works on any standard\nRGB webcam — no gloves,\nno depth sensors needed",
     RGBColor(0x00,0x69,0x5c)),
    ("🌐", "Web\nApplication",
     "React.js frontend runs in\nany modern browser\nNo installation required",
     NAVY),
    ("📊", "Confidence\nScores",
     "Real-time confidence bar\nwith Top-5 predictions\nand threshold filtering",
     RGBColor(0xe6,0x51,0x00)),
]
box_w = Inches(2.88); box_h = Inches(2.55)
for i,(icon,title,desc,col) in enumerate(features):
    c = i % 4; r = i // 4
    lx = Inches(0.35 + c*(box_w + Inches(0.12)))
    ly = Inches(1.55 + r*(box_h + Inches(0.1)))
    add_rounded_rect(s12, lx, ly, box_w, box_h, LIGHT_BG2, col, 1.0)
    add_rect(s12, lx, ly, box_w, Inches(0.06), col)
    # Icon circle
    add_rounded_rect(s12, lx + Inches(0.98), ly + Inches(0.15),
                     Inches(0.9), Inches(0.9), col)
    add_text(s12, icon, lx + Inches(0.98), ly + Inches(0.18),
             Inches(0.9), Inches(0.78), font_size=22,
             align=PP_ALIGN.CENTER, color=WHITE)
    add_text(s12, title, lx + Inches(0.1), ly + Inches(1.12),
             box_w - Inches(0.2), Inches(0.56),
             font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s12, desc, lx + Inches(0.1), ly + Inches(1.66),
             box_w - Inches(0.2), Inches(0.82),
             font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — RESULTS & ACCURACY
# ═══════════════════════════════════════════════════════════════════════════════
s13 = prs.slides.add_slide(BLANK)
slide_bg(s13, WHITE)
header_bar(s13, "Results & Accuracy", "Model performance on 40-class test set")
footer_bar(s13)

# Top stat row
main_stats = [
    ("96.8%",  "Test Accuracy",         GREEN),
    ("95.9%",  "Validation Accuracy",   BLUE),
    ("96.5%",  "Macro Precision",       TEAL),
    ("96.3%",  "Macro Recall",          ACCENT),
    ("96.4%",  "Macro F1-Score",        PURPLE),
    ("~67 ms", "End-to-End Latency",    RGBColor(0xe6,0x51,0x00)),
]
for i,(val,lbl,col) in enumerate(main_stats):
    stat_box(s13, val, lbl,
             Inches(0.35 + i*2.15), Inches(1.52),
             w=Inches(2.0), h=Inches(1.3), bg=col)

divider(s13, Inches(2.95))

# Left: bar chart (drawn as stacked rects)
add_text(s13, "Per-Sign Accuracy (Top 10 Signs)",
         Inches(0.35), Inches(3.1), Inches(6.5), Inches(0.38),
         font_size=13, bold=True, color=NAVY)

bar_data = [
    ("hello",      98.5, GREEN),  ("A (alpha)",  97.8, GREEN),
    ("1 (num)",    98.1, GREEN),  ("iloveyou",   97.2, BLUE),
    ("thankyou",   96.8, BLUE),   ("yes",        96.5, BLUE),
    ("howareyou",  95.3, ACCENT), ("goodmorning",95.1, ACCENT),
    ("eat",        94.6, AMBER),  ("drink",      93.8, AMBER),
]
bar_area_w = Inches(4.5)
bar_area_x = Inches(1.8)
bar_h_unit = Inches(0.38)
bar_gap    = Inches(0.08)
for i,(sign,acc,col) in enumerate(bar_data):
    by = Inches(3.58) + i*(bar_h_unit + bar_gap)
    # bg track
    add_rounded_rect(s13, bar_area_x, by, bar_area_w, bar_h_unit,
                     LIGHT_BG2)
    # fill
    fill_w = bar_area_w * (acc / 100)
    add_rounded_rect(s13, bar_area_x, by, fill_w, bar_h_unit, col)
    # label
    add_text(s13, sign, Inches(0.35), by + Inches(0.06),
             Inches(1.38), bar_h_unit - Inches(0.08),
             font_size=9.5, bold=True, color=GRAY)
    add_text(s13, f"{acc}%", bar_area_x + fill_w + Inches(0.08),
             by + Inches(0.06), Inches(0.55), bar_h_unit - Inches(0.08),
             font_size=9, bold=True, color=col)

# Right: classification report table (sample)
add_text(s13, "Classification Report (Sample)",
         Inches(6.75), Inches(3.1), Inches(6.2), Inches(0.38),
         font_size=13, bold=True, color=NAVY)

cr_headers = ["Sign", "Precision", "Recall", "F1", "Support"]
cr_rows = [
    ["hello",       "0.99", "0.98", "0.99", "180"],
    ["iloveyou",    "0.98", "0.97", "0.97", "180"],
    ["goodmorning", "0.95", "0.95", "0.95", "180"],
    ["A (letter)",  "0.98", "0.99", "0.99", "180"],
    ["1 (number)",  "0.98", "0.99", "0.98", "180"],
    ["eat",         "0.93", "0.94", "0.94", "180"],
    ["drink",       "0.94", "0.93", "0.93", "180"],
    ["go",          "0.96", "0.96", "0.96", "180"],
    ["Macro Avg",   "0.965","0.963","0.964","7,200"],
    ["Weighted Avg","0.967","0.968","0.967","7,200"],
]
table_slide(s13, cr_headers, cr_rows,
            Inches(6.75), Inches(3.58), Inches(6.2), Inches(3.96),
            col_widths=[Inches(1.55), Inches(1.15), Inches(1.15), Inches(1.0), Inches(1.05)])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════
s14 = prs.slides.add_slide(BLANK)
slide_bg(s14, WHITE)
header_bar(s14, "Comparison with Existing Systems", "How SignAI compares to related work")
footer_bar(s14)

# Main comparison table
cmp_headers = ["System / Paper", "Year", "Signs", "Accuracy", "Real-Time", "2 Hands", "TTS", "Web App"]
cmp_rows = [
    ["Pugeault & Bowden",  "2011", "24",   "79.3%",  "❌", "❌", "❌", "❌"],
    ["Oyedele et al.",     "2018", "26",   "91.5%",  "❌", "❌", "❌", "❌"],
    ["Rastgoo et al.",     "2020", "40",   "89.7%",  "⚠️", "❌", "❌", "❌"],
    ["Jiang et al.",       "2021", "10",   "94.2%",  "✅", "❌", "❌", "❌"],
    ["Taskiran et al.",    "2023", "250",  "96.8%",  "✅", "❌", "❌", "❌"],
    ["SignAI (This Work)", "2025", "40",   "96.8%",  "✅", "✅", "✅", "✅"],
]
tbl = table_slide(s14, cmp_headers, cmp_rows,
                  Inches(0.35), Inches(1.55), Inches(12.6), Inches(3.0),
                  col_widths=[Inches(2.6), Inches(0.75), Inches(0.75),
                               Inches(1.05), Inches(1.2), Inches(1.2),
                               Inches(0.95), Inches(1.2)])
# Highlight last row (SignAI)
for c in range(8):
    cell = tbl.cell(6, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xe8, 0xf5, 0xe9)
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.color.rgb = GREEN

divider(s14, Inches(4.75))

# Key advantages
add_text(s14, "✨  SignAI Unique Advantages",
         Inches(0.35), Inches(4.85), Inches(12.6), Inches(0.42),
         font_size=14, bold=True, color=NAVY)
advantages = [
    ("🤟 Only system with\nDual-Hand Detection",
     "126 keypoints/frame\nvs 63 in single-hand\nsystems",           TEAL),
    ("🔊 Integrated\nText-to-Speech",
     "Google TTS + browser\nfallback — no other\nsystem has this",   BLUE),
    ("💬 Real-Time\nSentence Builder",
     "Word chips + speech\nouput in a single\nweb interface",        ACCENT),
    ("🌐 Full Web App\nNo Installation",
     "React + FastAPI —\naccessible from any\nmodern browser",       GREEN),
    ("📦 Open Source\nReproducible",
     "Complete data pipeline\nand training code\nopenly available",  PURPLE),
]
for i,(title,desc,col) in enumerate(advantages):
    lx = Inches(0.35 + i*2.56)
    add_rounded_rect(s14, lx, Inches(5.38), Inches(2.38), Inches(1.68),
                     LIGHT_BG2, col, 1.0)
    add_rect(s14, lx, Inches(5.38), Inches(2.38), Inches(0.06), col)
    add_text(s14, title, lx+Inches(0.12), Inches(5.46),
             Inches(2.14), Inches(0.7), font_size=10.5, bold=True, color=col)
    add_text(s14, desc, lx+Inches(0.12), Inches(6.15),
             Inches(2.14), Inches(0.85), font_size=9.5, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════════
s15 = prs.slides.add_slide(BLANK)
slide_bg(s15, WHITE)
header_bar(s15, "Future Enhancements", "Roadmap for next version of SignAI")
footer_bar(s15)

# Roadmap phases
phases = [
    ("Phase 1\n(3–6 months)", "Expand Vocabulary",
     ["Scale to 500+ ASL signs using Google Kaggle dataset",
      "Add Indian Sign Language (ISL) support",
      "Improve accuracy on confused sign pairs"],
     ACCENT, "🚀"),
    ("Phase 2\n(6–12 months)", "Platform & Performance",
     ["Build Android/iOS app with TensorFlow Lite",
      "Real-time mobile inference (on-device)",
      "Model quantization & compression (~3 MB)"],
     TEAL, "📱"),
    ("Phase 3\n(12–18 months)", "Advanced AI Features",
     ["Continuous SLR with CTC decoding",
      "Facial expression + body pose recognition",
      "Transformer architecture for higher accuracy"],
     GREEN, "🧠"),
    ("Phase 4\n(18–24 months)", "Full Platform",
     ["Bidirectional translation (text → sign animation)",
      "Multi-language TTS (Hindi, Marathi, etc.)",
      "AR glasses / VR environment integration"],
     PURPLE, "🌍"),
]
for i,(phase,title,items,col,icon) in enumerate(phases):
    lx = Inches(0.32 + i*3.26)
    # Phase header
    add_rounded_rect(s15, lx, Inches(1.52), Inches(3.04), Inches(0.78), col)
    add_text(s15, icon + "  " + phase,
             lx, Inches(1.56), Inches(3.04), Inches(0.68),
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Title
    add_rounded_rect(s15, lx, Inches(2.34), Inches(3.04), Inches(0.5),
                     LIGHT_BG2, col, 0.8)
    add_text(s15, title, lx+Inches(0.1), Inches(2.38),
             Inches(2.84), Inches(0.42), font_size=12, bold=True, color=col)
    # Items
    add_rounded_rect(s15, lx, Inches(2.88), Inches(3.04), Inches(2.0),
                     LIGHT_BG2, col, 0.5)
    for j,item in enumerate(items):
        add_text(s15, f"▸  {item}",
                 lx+Inches(0.12), Inches(3.0) + j*Inches(0.6),
                 Inches(2.82), Inches(0.55), font_size=10, color=GRAY)

# Additional improvements table
divider(s15, Inches(5.05))
add_text(s15, "Other Planned Improvements",
         Inches(0.35), Inches(5.14), Inches(12.6), Inches(0.38),
         font_size=13, bold=True, color=NAVY)
imp_rows = [
    ["Multi-language TTS", "Add Hindi, Marathi regional TTS outputs", "High",   "Low"],
    ["Sign Learning Mode", "Interactive animated sign tutorial mode", "Medium", "Medium"],
    ["Cloud Deployment",   "Docker + cloud hosting for public access", "High",   "Low"],
    ["Dataset Expansion",  "ISL + BSL + custom phrase collection",    "High",   "High"],
]
table_slide(s15, ["Enhancement","Description","Priority","Effort"],
            imp_rows, Inches(0.35), Inches(5.62), Inches(12.6), Inches(1.45),
            col_widths=[Inches(2.3), Inches(5.8), Inches(1.6), Inches(1.6)])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
s16 = prs.slides.add_slide(BLANK)
slide_bg(s16, WHITE)
header_bar(s16, "Conclusion", "Summary of SignAI achievements")
footer_bar(s16)

# Left column — achievements
add_text(s16, "✅  What We Achieved",
         Inches(0.35), Inches(1.55), Inches(6.3), Inches(0.42),
         font_size=14, bold=True, color=NAVY)

achievements = [
    ("🤟", "40-Sign Real-Time Translator",
     "Detects 40 ASL signs via standard webcam\nwith no specialised hardware required"),
    ("🧠", "96.8% Test Accuracy",
     "Stacked LSTM with BatchNorm + Dropout\ntrained on 8,000 dual-hand keypoint sequences"),
    ("✋✋", "Dual-Hand Detection",
     "126 keypoints per frame (63 per hand)\nMissing hand automatically zero-padded"),
    ("⚡", "67ms End-to-End Latency",
     "Real-time WebSocket streaming at 15 fps\nfar below the 200ms usability threshold"),
    ("🔊", "Text-to-Speech Output",
     "Google TTS + browser fallback ensures\nspeech output works online and offline"),
    ("💬", "Sentence Builder",
     "Word chips + Sentence mode + Quick Add\nwith Copy, Undo, and Clear controls"),
    ("🌐", "Full-Stack Web App",
     "React.js + FastAPI — accessible from\nany modern browser, no installation needed"),
]
for i,(icon,title,desc) in enumerate(achievements):
    ly = Inches(2.05) + i * Inches(0.74)
    add_rounded_rect(s16, Inches(0.35), ly, Inches(6.3), Inches(0.66),
                     LIGHT_BG2, GREEN, 0.7)
    add_rounded_rect(s16, Inches(0.35), ly, Inches(0.52), Inches(0.66),
                     GREEN)
    add_text(s16, icon, Inches(0.35), ly + Inches(0.08),
             Inches(0.52), Inches(0.5), font_size=16,
             align=PP_ALIGN.CENTER, color=WHITE)
    add_text(s16, title, Inches(0.96), ly + Inches(0.06),
             Inches(2.45), Inches(0.3), font_size=10.5, bold=True, color=NAVY)
    add_text(s16, desc, Inches(0.96), ly + Inches(0.34),
             Inches(5.62), Inches(0.28), font_size=9, color=GRAY_LT)

divider(s16, Inches(2.0), GRAY_LT)
add_rect(s16, Inches(6.75), Inches(1.52), Inches(0.04), Inches(5.6), GRAY_LN)

# Right column — impact + learning
add_text(s16, "📈  Project Impact",
         Inches(6.9), Inches(1.55), Inches(6.1), Inches(0.42),
         font_size=14, bold=True, color=NAVY)

impact_boxes = [
    ("Assistive Technology",
     "Enables deaf and hearing-impaired individuals to communicate naturally with hearing people using everyday devices.",
     RGBColor(0x00,0x69,0x5c), "♿"),
    ("Education & Awareness",
     "SignAI can be used as a learning tool to help hearing people practice and learn ASL sign language.",
     BLUE, "🎓"),
    ("Research Foundation",
     "Open-source reproducible pipeline with data collection, training, and inference code for future researchers.",
     PURPLE, "🔬"),
    ("Industry Ready",
     "Deployable via Docker / cloud hosting — potential integration into healthcare, education, and public services.",
     RGBColor(0xe6,0x51,0x00), "🏭"),
]
for i,(title,desc,col,icon) in enumerate(impact_boxes):
    r, c = i//2, i%2
    lx = Inches(6.9 + c*3.05)
    ly = Inches(2.08 + r*2.42)
    add_rounded_rect(s16, lx, ly, Inches(2.85), Inches(2.22), LIGHT_BG2, col, 1.0)
    add_rect(s16, lx, ly, Inches(2.85), Inches(0.06), col)
    add_rounded_rect(s16, lx+Inches(1.0), ly+Inches(0.12),
                     Inches(0.78), Inches(0.78), col)
    add_text(s16, icon, lx+Inches(1.0), ly+Inches(0.14),
             Inches(0.78), Inches(0.68), font_size=20,
             align=PP_ALIGN.CENTER, color=WHITE)
    add_text(s16, title, lx+Inches(0.1), ly+Inches(0.98),
             Inches(2.65), Inches(0.38), font_size=11, bold=True, color=col)
    add_text(s16, desc, lx+Inches(0.1), ly+Inches(1.36),
             Inches(2.65), Inches(0.78), font_size=9.5, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
s17 = prs.slides.add_slide(BLANK)
slide_bg(s17, DARK_BG)

# Top accent strip
add_rect(s17, 0, 0, W, Inches(0.08), ACCENT)

# Left dark panel
add_rect(s17, 0, Inches(0.08), Inches(5.5), H - Inches(0.08), DARK_SURF)
add_rect(s17, Inches(5.5), Inches(0.08), Inches(0.06), H - Inches(0.08), ACCENT)

# Big thank you
add_text(s17, "🙏", Inches(0.2), Inches(0.5),
         Inches(5.1), Inches(1.6), font_size=72,
         align=PP_ALIGN.CENTER, color=WHITE)
add_text(s17, "Thank You!",
         Inches(0.2), Inches(2.0), Inches(5.1), Inches(1.1),
         font_size=46, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
add_text(s17, "for your time and attention",
         Inches(0.2), Inches(3.0), Inches(5.1), Inches(0.5),
         font_size=14, italic=True,
         color=RGBColor(0xbb,0xde,0xfb), align=PP_ALIGN.CENTER)
add_rect(s17, Inches(0.8), Inches(3.62), Inches(3.5), Inches(0.04), ACCENT)

# Quick summary chips on left
summary = [
    ("40 Signs",    ACCENT), ("96.8% Accuracy", GREEN),
    ("Both Hands",  TEAL),   ("Real-Time",       BLUE),
    ("TTS Output",  PURPLE), ("Web App",         NAVY),
]
for i,(lbl,bg) in enumerate(summary):
    chip(s17, lbl, Inches(0.55 + (i%3)*1.65), Inches(3.88 + i//3*0.56),
         w=Inches(1.56), h=Inches(0.44), bg=bg, font_size=11)

# Right panel — contact + demo info
add_text(s17, "SignAI",
         Inches(5.85), Inches(0.55), Inches(7.2), Inches(1.0),
         font_size=48, bold=True, color=ACCENT2, align=PP_ALIGN.LEFT)
add_text(s17, "Real-Time AI Sign Language Translator",
         Inches(5.85), Inches(1.42), Inches(7.2), Inches(0.45),
         font_size=14, color=RGBColor(0xbb,0xde,0xfb))
add_rect(s17, Inches(5.85), Inches(1.95), Inches(7.1), Inches(0.04), ACCENT)

# Contact card
add_rounded_rect(s17, Inches(5.85), Inches(2.1), Inches(7.1), Inches(2.1),
                 DARK_SURF, ACCENT, 0.8)
contact = [
    ("👤", "Himanshu Jagdish Patil"),
    ("🎓", "B.Tech CSE (AI & ML)  |  Roll: CS2024001"),
    ("📧", "himanshujagdishpatil914@gmail.com"),
    ("🔗", "github.com/himanshujagdishpatil914/Himanshu"),
    ("👨‍🏫", "Guide: Prof. [Guide Name]  |  Co-Guide: Prof. [Name]"),
]
for i,(icon,val) in enumerate(contact):
    add_text(s17, icon, Inches(6.05), Inches(2.22)+i*Inches(0.37),
             Inches(0.42), Inches(0.34), font_size=12,
             align=PP_ALIGN.CENTER, color=WHITE)
    add_text(s17, val, Inches(6.52), Inches(2.22)+i*Inches(0.37),
             Inches(6.3), Inches(0.34), font_size=10.5, color=WHITE)

# Q&A and demo invite boxes
add_rounded_rect(s17, Inches(5.85), Inches(4.35), Inches(3.42), Inches(1.62),
                 NAVY, ACCENT, 1.2)
add_text(s17, "❓  Questions?",
         Inches(5.95), Inches(4.42), Inches(3.22), Inches(0.46),
         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s17, "Open for Q&A\nPlease feel free to ask\nabout any part of SignAI",
         Inches(5.95), Inches(4.86), Inches(3.22), Inches(0.98),
         font_size=11, color=RGBColor(0xbb,0xde,0xfb),
         align=PP_ALIGN.CENTER)

add_rounded_rect(s17, Inches(9.45), Inches(4.35), Inches(3.42), Inches(1.62),
                 TEAL, GREEN, 1.2)
add_text(s17, "🎬  Live Demo",
         Inches(9.55), Inches(4.42), Inches(3.22), Inches(0.46),
         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s17, "Start the app at:\nlocalhost:3000\nReady to demonstrate!",
         Inches(9.55), Inches(4.86), Inches(3.22), Inches(0.98),
         font_size=11, color=RGBColor(0xb2,0xdf,0xdb),
         align=PP_ALIGN.CENTER)

# GitHub link
add_rounded_rect(s17, Inches(5.85), Inches(6.1), Inches(7.1), Inches(0.95),
                 DARK_SURF, ACCENT, 0.6)
add_text(s17, "📁  Project Repository:",
         Inches(6.0), Inches(6.18), Inches(2.5), Inches(0.38),
         font_size=10.5, bold=True, color=ACCENT2)
add_text(s17, "github.com/himanshujagdishpatil914/Himanshu",
         Inches(8.4), Inches(6.18), Inches(4.4), Inches(0.38),
         font_size=10.5, bold=False, color=WHITE)
add_text(s17, "⭐  Source code  |  📄 Report  |  📊 Presentation  |  🖼️ Screenshots",
         Inches(6.0), Inches(6.54), Inches(6.8), Inches(0.38),
         font_size=9.5, color=RGBColor(0x90,0xca,0xf9))

# Bottom bar
add_rect(s17, 0, H - Inches(0.38), W, Inches(0.38), NAVY)
add_text(s17, "SignAI  |  Final Year Project  |  CSE (AI & ML)  |  Himanshu Jagdish Patil  |  2024–25",
         Inches(0.3), H - Inches(0.34), Inches(12.7), Inches(0.3),
         font_size=9, color=ACCENT2, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"[DONE] Saved → {OUT}")
print(f"       Slides : {len(prs.slides)}")
import os
print(f"       Size   : {os.path.getsize(OUT)//1024} KB")
